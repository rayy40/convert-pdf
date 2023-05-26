from fastapi import FastAPI, Request
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Union, List, Dict
from PyPDF2 import PdfReader, PdfWriter
from dotenv import load_dotenv
from pdf2docx import Converter
from pptx import Presentation
import aspose.words as aw
from io import BytesIO
from PIL import Image
import urllib.request
import requests
import zipfile
import pyrebase
import math
import fitz
import os

load_dotenv()

firebase_config = {
    "apiKey": os.getenv("API_KEY"),
    "authDomain": os.getenv("AUTH_DOMAIN"),
    "projectId": os.getenv("PROJECT_ID"),
    "storageBucket": os.getenv("STORAGE_BUCKET"),
    "messagingSenderId": os.getenv("MESSAGING_SENDER_ID"),
    "appId": os.getenv("APP_ID"),
    "measurementId": os.getenv("MEASUREMENT_ID"),
    "databaseURL": "",
}

app = FastAPI()

# Enable CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class PDFManipulationRequest(BaseModel):
    urls: Union[str, List[str]]
    pages: List[int] = []
    rotation: Dict = {}
    password: str = ""
    metadata: Dict


@app.post("/api/jpg-to-pdf")
async def jpg_to_pdf(request: Request, data: PDFManipulationRequest):
    urls = data.urls
    metadata = data.metadata
    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    # Create a new PDF document
    pdf = fitz.open()

    for url in urls:
        # Download image from URL
        with urllib.request.urlopen(url) as img_url:
            img_data = img_url.read()

        # Load the image data into a PIL Image object
        img = Image.open(BytesIO(img_data))

        a4_width = 595.0  # A4 page width
        a4_height = 842.0  # A4 page height
        margin = 28.35  # 1 cm
        max_width = a4_width - 2 * margin
        max_height = a4_height - 2 * margin

        # Resize the image while maintaining its aspect ratio
        img.thumbnail((max_width, max_height))

        # Add new page with A4 dimensions
        page = pdf.new_page(width=a4_width, height=a4_height)

        # Get the dimensions of the resized image
        img_width, img_height = img.size

        # Calculate position to center image
        x_pos = (a4_width - img_width) / 2
        y_pos = (a4_height - img_height) / 2

        # Convert the PIL Image to a Fitz Pixmap
        pixmap = fitz.Pixmap(BytesIO(img_data))

        # Insert the image into the PDF page
        page.insert_image(fitz.Rect(x_pos, y_pos, x_pos + img_width, y_pos + img_height), pixmap=pixmap)

    # Save PDF to file and return download URL
    pdf_data = BytesIO()
    pdf.save(pdf_data)
    pdf.close()

    # Set the destination path for the file upload
    folder_name = "jpg-to-pdf/modified"
    destination_path = f"{folder_name}/{filename_without_extension}.pdf"

    # Upload the single image file to Firebase Storage
    firebase = pyrebase.initialize_app(firebase_config)
    storage = firebase.storage()
    storage.child(destination_path).put(pdf_data.getvalue())

    # Get the download URL
    download_url = storage.child(destination_path).get_url(None)
    
    pdf_data.close()

    return download_url


@app.post("/api/pdf-to-jpg")
async def pdf_to_jpg(request: Request, data: PDFManipulationRequest):
    urls = data.urls
    metadata = data.metadata
    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    # Create a BytesIO object to store the images
    images_data = BytesIO()

    # Download the PDF file
    response = requests.get(urls[0])
    pdf_data = BytesIO(response.content)

    # Open the PDF
    with fitz.open(stream=pdf_data.read(), filetype="pdf") as pdf:
        # Calculate the number of pages
        num_pages = pdf.page_count

        # Iterate over each page
        for i, page in enumerate(pdf):
            # Get the page as a pixmap
            # page = pdf.load_page(i)
            pix = page.get_pixmap()

            # Save the pixmap as an image file
            image_data = pix.tobytes("jpeg")
            images_data.write(image_data)

    # Check the number of images
    if num_pages > 1:
        # Set the destination path for the file upload
        folder_name = "pdf-to-jpg/modified"
        destination_path = f"{folder_name}/{filename_without_extension}.zip"
        
        zip_stream = BytesIO()
        with zipfile.ZipFile(zip_stream, "w", zipfile.ZIP_STORED) as zipf:
            # Add each image to the zip file
            for i in range(num_pages):
                image_data = images_data.getvalue()
                zipf.writestr(f"{filename_without_extension}_{i}.jpg", image_data)

        # Reset the position of the zip_data object to the beginning
        zip_stream.seek(0)

        # Upload the zip file to Firebase Storage
        firebase = pyrebase.initialize_app(firebase_config)
        storage = firebase.storage()
        storage.child(destination_path).put(zip_stream.getvalue())

        # Generate the download URL and metadata for the zip file
        zip_url = storage.child(destination_path).get_url(None)

        # Clean up
        pdf_data.close()
        zip_stream.close()
        images_data.close()

        return zip_url
    elif num_pages == 1:
        # Set the destination path for the file upload
        folder_name = "pdf-to-jpg/modified"
        destination_path = f"{folder_name}/{filename_without_extension}.jpg"

        # Upload the single image file to Firebase Storage
        firebase = pyrebase.initialize_app(firebase_config)
        storage = firebase.storage()
        storage.child(destination_path).put(images_data.getvalue())

        # Generate the download URL for the single image
        image_url = storage.child(destination_path).get_url(None)

        # Clean up
        pdf_data.close()
        images_data.close()

        return image_url


@app.post("/api/pdf-to-word")
async def pdf_to_word(request: Request, data: PDFManipulationRequest):
    urls = data.urls
    metadata = data.metadata
    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    # Download the PDF file
    response = requests.get(urls[0])
    pdf_stream = BytesIO(response.content)

    cv = Converter(pdf_stream.read())

    output_doc_stream = BytesIO()
    cv.convert(output_doc_stream, start=0, end=None)
    cv.close()

    # Set the destination path for the file upload
    folder_name = "pdf-to-word/modified"
    destination_path = f"{folder_name}/{filename_without_extension}.docx"

    # Upload the single image file to Firebase Storage
    firebase = pyrebase.initialize_app(firebase_config)
    storage = firebase.storage()
    storage.child(destination_path).put(output_doc_stream.getvalue())

    # Generate the download URL for the modified PDF
    docx_url = storage.child(destination_path).get_url(None)

    # Clean up temporary files
    pdf_stream.close()
    output_doc_stream.close()

    return docx_url


@app.post("/api/word-to-pdf")
async def word_to_pdf(request: Request, data: PDFManipulationRequest):
    urls = data.urls
    metadata = data.metadata
    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    # Download the PDF file
    response = requests.get(urls[0])
    doc_stream = BytesIO(response.content)

    doc = aw.Document(stream=doc_stream)

    output_pdf_stream = BytesIO()
    doc.save(output_pdf_stream, aw.SaveFormat.PDF)

    # Set the destination path for the file upload
    folder_name = "word-to-pdf/modified"
    destination_path = f"{folder_name}/{filename_without_extension}.pdf"

    # Upload the single image file to Firebase Storage
    firebase = pyrebase.initialize_app(firebase_config)
    storage = firebase.storage()
    storage.child(destination_path).put(output_pdf_stream.getvalue())

    # Generate the download URL for the modified PDF
    docx_url = storage.child(destination_path).get_url(None)

    doc_stream.close()
    output_pdf_stream.close()

    return docx_url


@app.post("/api/delete-pages")
async def delete_pages(request: Request, data: PDFManipulationRequest):
    urls = data.urls
    pages = data.pages
    metadata = data.metadata
    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    response = requests.get(urls)
    pdf_stream = BytesIO(response.content)

    pdf = fitz.open(stream=pdf_stream.read(), filetype="pdf")

    for page_number in sorted(pages, reverse=True):
        pdf.delete_page(page_number - 1)

    output_pdf_stream = BytesIO()
    pdf.save(output_pdf_stream)
    pdf.close()

    # Set the destination path for the file upload
    folder_name = "delete-pages/modified"
    destination_path = f"{folder_name}/{filename_without_extension}-modified.pdf"

    # Upload the single image file to Firebase Storage
    firebase = pyrebase.initialize_app(firebase_config)
    storage = firebase.storage()
    storage.child(destination_path).put(output_pdf_stream.getvalue())

    # Generate the download URL for the modified PDF
    pdf_url = storage.child(destination_path).get_url(None)
    
    output_pdf_stream.close()
    pdf_stream.close()

    return pdf_url


@app.post("/api/rotate-pdf")
def rotate_pdf(request: Request, data: PDFManipulationRequest):
    urls = data.urls
    pages = data.pages
    rotation = data.rotation
    metadata = data.metadata
    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    response = requests.get(urls)
    pdf_stream = BytesIO(response.content)

    pdf = fitz.open(stream=pdf_stream.read(), filetype="pdf")

    for page_number in sorted(pages, reverse=True):
        page = pdf[page_number - 1]
        page.set_rotation(rotation[str(page_number)])

    output_pdf_stream = BytesIO()
    pdf.save(output_pdf_stream)
    pdf.close()

    # Set the destination path for the file upload
    folder_name = "rotate-pdf/modified"
    destination_path = f"{folder_name}/{filename_without_extension}-rotated.pdf"

    # Upload the single image file to Firebase Storage
    firebase = pyrebase.initialize_app(firebase_config)
    storage = firebase.storage()
    storage.child(destination_path).put(output_pdf_stream.getvalue())

    # Generate the download URL for the modified PDF
    pdf_url = storage.child(destination_path).get_url(None)
    
    pdf_stream.close()
    output_pdf_stream.close()

    return pdf_url


@app.post("/api/extract-pdf")
def extract_pdf(request: Request, data: PDFManipulationRequest):
    urls = data.urls
    pages = data.pages
    metadata = data.metadata
    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    response = requests.get(urls, stream=True)

   # Create in-memory file-like object to store the PDF content
    pdf_stream = BytesIO()
    
    for chunk in response.iter_content(chunk_size=4096):
        pdf_stream.write(chunk)

    # Extract selected pages
    doc = fitz.open(stream=pdf_stream, filetype="pdf")
    pdf = fitz.open()

    for page_num in sorted(pages, reverse=True): 
        pdf.insert_pdf(doc, from_page=int(page_num) - 1, to_page=int(page_num) - 1)

    # Save the extracted pages to a new PDF file
    output_pdf_stream = BytesIO()
    pdf.save(output_pdf_stream)
    pdf.close()
    doc.close()

    # Set the destination path for the file upload
    folder_name = "extract-pdf/modified"
    destination_path = f"{folder_name}/{filename_without_extension}-extracted.pdf"

    # Upload the new PDF file to Firebase storage
    firebase = pyrebase.initialize_app(firebase_config)
    storage = firebase.storage()
    storage.child(destination_path).put(output_pdf_stream.getvalue())

    # Get the download URL
    download_url = storage.child(destination_path).get_url(None)

    # Remove the local files
    pdf_stream.close()
    output_pdf_stream.close()

    return download_url


@app.post("/api/extract-images")
def extract_images(request: Request, data: PDFManipulationRequest):
    urls = data.urls
    metadata = data.metadata
    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    response = requests.get(urls[0])
    
    pdf_stream = BytesIO(response.content)
    pdf_document = fitz.open(stream=pdf_stream.read(), filetype="pdf")

    image_data = []

    for current_page in range(pdf_document.page_count):
        page = pdf_document.load_page(current_page)
        image_list = page.get_images()
        for image_index, image in enumerate(image_list):
            xref = image[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_extension = base_image["ext"]
            image_name = f"image{current_page}_{image_index}.{image_extension}"
            image_data.append((image_name, image_bytes))

    # Create a BytesIO object to store the zip file content
    zip_stream = BytesIO()

    # Create a zip file of the extracted images
    with zipfile.ZipFile(zip_stream, "w", zipfile.ZIP_STORED) as zip_file:
        for image_name, image_bytes in image_data:
            zip_file.writestr(image_name, image_bytes)
            
    # Set the position to the beginning of the BytesIO object
    zip_stream.seek(0)

    # Set the destination path for the file upload
    zip_file_name = f"{filename_without_extension}_extracted-images.zip"
    folder_name = "extract-images/modified"
    destination_path = f"{folder_name}/{zip_file_name}"

    # Upload the new PDF file to Firebase storage
    firebase = pyrebase.initialize_app(firebase_config)
    storage = firebase.storage()
    storage.child(destination_path).put(zip_stream.getvalue())

    # Get the download URL
    download_url = storage.child(destination_path).get_url(None)

    zip_stream.close()
    pdf_stream.close()
    pdf_document.close()

    return download_url


@app.post("/api/pdf-to-ppt")
def pdf_to_ppt(request: Request, data: PDFManipulationRequest):
    # Get the JSON data containing URLs
    urls = data.urls
    metadata = data.metadata
    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    response = requests.get(urls[0])

    # Load the PDF content into a Presentation object
    prs = Presentation()

    # Create a slide for each page in the PDF
    pdf_stream = BytesIO(response.content)
    pdf_stream.seek(0)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide_layout = prs.slide_layouts[1]  # Use the layout for a content slide


    pdf_doc = fitz.open(stream=pdf_stream.read(), filetype="pdf")
    for page_index in range(len(pdf_doc)):
        slide = prs.slides.add_slide(slide_layout)
        content_slide = slide.shapes

        # Convert PDF page to image
        page = pdf_doc.load_page(page_index)
        pix = page.get_pixmap()

        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

        image_stream = BytesIO()
        img.save(image_stream, format="PNG")
        image_stream.seek(0)

        # Calculate the dimensions while maintaining the aspect ratio
        max_width = slide_width
        max_height = slide_height
        img_width, img_height = img.size
        aspect_ratio = img_width / img_height

        if aspect_ratio > max_width / max_height:
            width = max_width
            height = math.ceil(width / aspect_ratio)
            left = 0
            top = (max_height - height) / 2
        else:
            height = max_height
            width = math.ceil(height * aspect_ratio)
            top = 0
            left = (max_width - width) / 2

        # Add the image to the slide
        content_slide.add_picture(image_stream, left, top, width, height)

    # Save the Presentation as a PPTX file
    converted_pptx_stream = BytesIO()
    prs.save(converted_pptx_stream)
    
    # Set the position to the beginning of the BytesIO object
    converted_pptx_stream.seek(0)

    # Set the destination path for the file upload
    folder_name = "pdf-to-ppt/modified"
    destination_path = f"{folder_name}/{filename_without_extension}.pptx"

    # Upload the converted PPTX file to Firebase Storage
    firebase = pyrebase.initialize_app(firebase_config)
    storage = firebase.storage()
    storage.child(destination_path).put(converted_pptx_stream.getvalue())

    # Generate the download URL for the converted PPTX
    pptx_url = storage.child(destination_path).get_url(None)
    
    pdf_stream.close()
    pdf_doc.close()
    converted_pptx_stream.close()

    return pptx_url


@app.post("/api/compress-pdf")
def compress_pdf(request: Request, data: PDFManipulationRequest):
    urls = data.urls
    metadata = data.metadata
    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    response = requests.get(urls[0])
    
    # Create a BytesIO object to store the PDF content
    pdf_stream = BytesIO(response.content)

    renderer = aw.pdf2word.fixedformats.PdfFixedRenderer()
    pdf_read_options = aw.pdf2word.fixedformats.PdfFixedOptions()
    pdf_read_options.image_format = aw.pdf2word.fixedformats.FixedImageFormat.JPEG
    pdf_read_options.jpeg_quality = 50

    pages_stream = renderer.save_pdf_as_images(pdf_stream, pdf_read_options)

    builder = aw.DocumentBuilder()
    for i in range(0, len(pages_stream)):
        # Set maximum page size to avoid the current page image scaling.
        max_page_dimension = 1584
        page_setup = builder.page_setup
        page_setup.page_width = max_page_dimension
        page_setup.page_height = max_page_dimension

        page_image = builder.insert_image(pages_stream[i])

        page_setup.page_width = page_image.width
        page_setup.page_height = page_image.height
        page_setup.top_margin = 0
        page_setup.left_margin = 0
        page_setup.bottom_margin = 0
        page_setup.right_margin = 0

        if i != len(pages_stream) - 1:
            builder.insert_break(aw.BreakType.SECTION_BREAK_NEW_PAGE)

    save_options = aw.saving.PdfSaveOptions()
    save_options.cache_background_graphics = True

    # Create a BytesIO object to store the compressed PDF content
    compressed_pdf_stream = BytesIO()
        
    builder.document.save(compressed_pdf_stream, save_options)

    # Set the destination path for the file upload
    folder_name = "compress-pdf/modified"
    destination_path = f"{folder_name}/{filename_without_extension}-compressed.pdf"

    # Upload the single image file to Firebase Storage
    firebase = pyrebase.initialize_app(firebase_config)
    storage = firebase.storage()
    storage.child(destination_path).put(compressed_pdf_stream.getvalue())

    # Generate the download URL for the modified PDF
    docx_url = storage.child(destination_path).get_url(None)
    
    pdf_stream.close()
    compressed_pdf_stream.close()

    return docx_url


@app.post("/api/protect-pdf")
def protect_pdf(request: Request, data: PDFManipulationRequest):
    urls = data.urls
    password = data.password
    metadata = data.metadata

    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    response = requests.get(urls)
    
    # Create a BytesIO object to store the PDF content
    pdf_stream = BytesIO(response.content)

    # Open the PDF
    pdf = PdfReader(pdf_stream)

    # Create a new PDF with encryption
    writer = PdfWriter()

    # Iterate through the pages of the original PDF and add them to the new PDF
    for page in pdf.pages:
        writer.add_page(page)

    # Encrypt the new PDF with a password
    writer.encrypt(user_password=password, owner_password=None, use_128bit=True)
    
    # Create a BytesIO object to store the encrypted PDF content
    encrypted_pdf_stream = BytesIO()

    # Save the encrypted PDF to the temporary file
    writer.write(encrypted_pdf_stream)

    # Set the position to the beginning of the BytesIO object
    encrypted_pdf_stream.seek(0)

    # Set the destination path for the file upload
    folder_name = "protect-pdf/modified"
    destination_path = f"{folder_name}/{filename_without_extension}-protected.pdf"

    # Upload the new PDF file to Firebase storage
    firebase = pyrebase.initialize_app(firebase_config)
    storage = firebase.storage()
    storage.child(destination_path).put(encrypted_pdf_stream.getvalue())

    # Get the download URL
    download_url = storage.child(destination_path).get_url(None)
    
    pdf_stream.close()
    encrypted_pdf_stream.close()

    return download_url


# @app.route("/api/unlock-pdf", methods=["POST"])
# @cross_origin(supports_credentials=True)
# def unlock_pdf():
#     return "Working on"


@app.post("/api/merge-pdf")
def merge_pdf(request: Request, data: PDFManipulationRequest):
    urls = data.urls
    metadata = data.metadata

    filename = metadata["name"].split("--")[0]
    filename_without_extension = os.path.splitext(filename)[0]

    merged_pdf = fitz.open()

    for url in urls:
        # Download the PDF file from the URL
        response = requests.get(url)
        pdf_stream = BytesIO(response.content)
        with fitz.open(stream=pdf_stream, filetype="pdf") as pdf_file:
            merged_pdf.insert_pdf(pdf_file)
            
    output_stream = BytesIO()
    merged_pdf.save(output_stream)
    merged_pdf.close()

    # Set the destination path for the file upload
    folder_name = "merge-pdf/modified"
    destination_path = f"{folder_name}/{filename_without_extension}-merged.pdf"

    # Upload the single image file to Firebase Storage
    firebase = pyrebase.initialize_app(firebase_config)
    storage = firebase.storage()
    storage.child(destination_path).put(output_stream.getvalue())

    # Generate the download URL for the modified PDF
    pdf_url = storage.child(destination_path).get_url(None)
    
    output_stream.close()
    pdf_stream.close()

    return pdf_url

